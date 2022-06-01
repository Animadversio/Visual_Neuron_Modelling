import matplotlib.pyplot as plt
import pptx
from pptx import Presentation, slide
from pptx.util import Inches
from pptx_utils import *
from scipy.io import loadmat
mat_path = r"E:\OneDrive - Washington University in St. Louis\Mat_Statistics"
#%%
rootdir = r"E:\OneDrive - Harvard University\Manifold_NeuralRegress"
sumdir = r"E:\OneDrive - Harvard University\Manifold_NeuralRegress\summary"
figdir = join(sumdir, "per_experiment")
view_layout_params(join(sumdir, "ManifExp_all.pptx"), 1)
#%%
def chan2area(chan):
    if chan <= 32 and chan >= 1:
        return "IT"
    elif chan >= 49 and chan <=64:
        return "V4"
    elif chan > 32 and chan < 49:
        return "V1"
    else:
        raise ValueError("Invalid channel number")


def layout_slide(slide, title_text, protopath, figpath):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.top = Inches(0.0)
    tf.left = Inches(0.0)
    tf.width = Inches(6.66667)
    tf.height = Inches(0.83333)
    tf.text_frame._set_font("Candara", 48, False, False)
    pic = slide.shapes.add_picture(protopath, Inches(0.0), Inches(0.83333), width=Inches(6.66667))
    pic = slide.shapes.add_picture(figpath, Inches(7.7337), Inches(0.0), height=Inches(7.5))


def layout_first_slide(slide, title_text, protopath, figpath):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.top = Inches(0.0)
    tf.left = Inches(0.0)
    tf.width = Inches(4.833333)
    tf.height = Inches(1.6145833)
    tf.text_frame._set_font("Candara", 36, False, False)
    pic0 = slide.shapes.add_picture(protopath, Inches(0.9166666), Inches(2.541666), width=Inches(2.666667))
    pic = slide.shapes.add_picture(figpath, Inches(4.59449), Inches(0.0), height=Inches(7.5), width=Inches(8.7388396))
    picdict = {'crop_right': 0.09167, 'crop_left': 0.09271, 'crop_top': 0.0, 'crop_bottom': 0.06666}
    for k, v in picdict.items():
        setattr(pic, k, v)

#%%
prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[5]
for Animal in ["Alfa", "Beto"]:
    for Expi in range(1, 47):
        if Animal == "Beto" and Expi == 46: continue
        for featlayer in [".layer2.Bottleneck3", ".layer3.Bottleneck5", ".layer4.Bottleneck2"]:
            protopath = join(figdir, f"{Animal}-Exp{Expi:02d}-{featlayer}-regr_merge_vis.png")
            barfigpath = join(figdir, f"prediction_comparison_{Animal}_{Expi:02d}_{featlayer}.png")
            title_text = f"{Animal} Exp{Expi:02d} {featlayer.split('.')[1]}"
            slide = prs.slides.add_slide(blank_slide_layout)
            layout_slide(slide, title_text, protopath, barfigpath)
#%%
prs.save(join(sumdir, "ManifExp_regression_pptx_export.pptx"))
#%%
MStats = loadmat(join(mat_path, Animal + "_Manif_stats.mat"), struct_as_record=False, squeeze_me=True)['Stats']
EStats = loadmat(join(mat_path, Animal + "_Evol_stats.mat"), struct_as_record=False, squeeze_me=True, chars_as_strings=True)['EStats']
ReprStats = loadmat(join(mat_path, Animal + "_ImageRepr.mat"), struct_as_record=False, squeeze_me=True, chars_as_strings=True)['ReprStats']
protoimg = ReprStats[Expi-1].Manif.BestImg

#%%

#%%
corrfeat_figdir = r"E:\OneDrive - Washington University in St. Louis\corrFeatTsr_FactorVis\models\resnet50_linf8-layer3_NF3_bdr1_Tthresh_3__nobdr_res-robust_CV"

prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[5]
for Animal in ["Alfa", "Beto"]:
    EStats = loadmat(join(mat_path, Animal + "_Evol_stats.mat"),
        struct_as_record=False, squeeze_me=True, chars_as_strings=True)['EStats']
    for Expi in range(1, 47):
        if Animal == "Beto" and Expi == 46: continue
        pref_chan = EStats[Expi-1].evol.pref_chan
        imgpos = EStats[Expi-1].evol.imgpos
        imgsize = EStats[Expi-1].evol.imgsize
        area = chan2area(pref_chan)
        title_str = f"{Animal} Exp{Expi:02d} {area} Pref Ch{pref_chan}\npos: {imgpos} {imgsize} deg"
        ccfigpath = join(corrfeat_figdir, f"{Animal}_Exp{Expi:02d}_summary.png")
        protopath = join(sumdir, "proto", f"{Animal}_Exp{Expi:02d}_manif_proto.png")
        slide = prs.slides.add_slide(blank_slide_layout)
        layout_first_slide(slide, title_str, protopath, ccfigpath)
        for featlayer in [".layer2.Bottleneck3", ".layer3.Bottleneck5", ".layer4.Bottleneck2"]:
            protopath = join(figdir, f"{Animal}-Exp{Expi:02d}-{featlayer}-regr_merge_vis.png")
            barfigpath = join(figdir, f"prediction_comparison_{Animal}_{Expi:02d}_{featlayer}.png")
            title_text = f"{Animal} Exp{Expi:02d} {featlayer.split('.')[1]}"
            slide = prs.slides.add_slide(blank_slide_layout)
            layout_slide(slide, title_text, protopath, barfigpath)

prs.save(join(sumdir, "ManifExp_regression_pptx_cc_export.pptx"))
#%%
#%%
def layout_proto_evol_slide(slide, title_text, protopath, evol_figpath, manif_figpath):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.text_frame._set_font("Candara", 36, False, False)
    for k, v in {'height': 1.690, 'width': 9.313, 'top': 0.0, 'left': 3.103}.items():
        setattr(tf, k, Inches(v))
    pic0 = slide.shapes.add_picture(protopath, Inches(0.0), Inches(0.0), width=Inches(2.60))
    pic1 = slide.shapes.add_picture(evol_figpath, Inches(0.0), Inches(0.0), )
    pic2 = slide.shapes.add_picture(manif_figpath, Inches(0.0), Inches(0.0), )
    for k, v in {'height': 5.600, 'width': 6.109, 'top': 1.90, 'left': 1.624}.items():
        setattr(pic1, k, Inches(v))
    for k, v in {'height': 5.600, 'width': 5.600, 'top': 1.90, 'left': 7.733}.items():
        setattr(pic2, k, Inches(v))


def layout_slide_factor(slide, title_text, protomtgpath, figpath):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.text_frame._set_font("Candara", 36, False, False)
    pic1 = slide.shapes.add_picture(protomtgpath, Inches(0.0), Inches(0.83333), )
    pic2 = slide.shapes.add_picture(figpath, Inches(7.7337), Inches(0.0), )
    for k, v in {'height': 1.767, 'width': 2.191, 'top': 4.951, 'left': 0.0}.items():
        setattr(tf, k, Inches(v))
    for k,v in {'height': 2.739, 'width': 10.924, 'top': 4.761, 'left': 2.409}.items():
        setattr(pic1, k, Inches(v))
    for k,v in {'height': 4.761, 'width': 11.796, 'top': 0.0, 'left': 1.537}.items():
        setattr(pic2, k, Inches(v))
#%%
layoutdict = view_layout_params(join(sumdir, "ManifExp_all.pptx"), 4)
#%%  Export Evolution trajectory and Manifold images?
"""
Slide design: 4 slides per experiment
- one slide for evolution traces and manifold map 
- three slides for each feature layer: 
    - summarizing bar plot 
    - proto montage. 
"""
prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[5]
for Animal in ["Alfa", "Beto"]: #
    EStats = loadmat(join(mat_path, Animal + "_Evol_stats.mat"),
        struct_as_record=False, squeeze_me=True, chars_as_strings=True)['EStats']
    for Expi in tqdm(range(1, 47)):
        if Animal == "Beto" and Expi == 46: continue
        pref_chan = EStats[Expi-1].evol.pref_chan
        imgpos = EStats[Expi-1].evol.imgpos
        imgsize = EStats[Expi-1].evol.imgsize
        area = chan2area(pref_chan)
        title_str = f"{Animal} Exp{Expi:02d} {area} Pref Ch{pref_chan}\npos: {imgpos} {imgsize} deg"
        protopath = join(sumdir, "proto", f"{Animal}_Exp{Expi:02d}_manif_proto.png")
        evolfigpath = join(sumdir, "classic_figs", f"{Animal}_Exp{Expi:02d}_Evolution.png")
        maniffigpath = join(sumdir, "classic_figs", f"{Animal}_Exp{Expi:02d}_Manifold.png")
        # ccfigpath = join(corrfeat_figdir, f"{Animal}_Exp{Expi:02d}_summary.png")
        slide = prs.slides.add_slide(blank_slide_layout)
        layout_proto_evol_slide(slide, title_str, protopath, evolfigpath, maniffigpath)
        for featlayer in [".layer2.Bottleneck3", ".layer3.Bottleneck5", ".layer4.Bottleneck2"]:
            protopath = join(figdir, f"{Animal}-Exp{Expi:02d}-{featlayer}-factregr_merge_vis.png")
            barfigpath = join(figdir, f"prediction_comparison_factor_Trsp_{Animal}_{Expi:02d}_{featlayer}.png")
            title_text = f"{Animal} Exp{Expi:02d} {featlayer.split('.')[1]}"
            slide = prs.slides.add_slide(blank_slide_layout)
            layout_slide_factor(slide, title_text, protopath, barfigpath)

prs.save(join(sumdir, "ManifExp_factor_regress_pptx_export.pptx"))
#%%
import pandas as pd
def dataframe2pptxtable(df, graphtab, font=("Candara", 12, False, False)):
    """
    Convert a pandas dataframe to a pptx table.
    """
    for c in range(df.shape[1]):
        graphtab.table.cell(0, c).text = df.columns[c]
        graphtab.table.cell(0, c).text_frame._set_font(*font)
        for r in range(df.shape[0]):
            entry = df.iloc[r, c]
            if np.issubdtype(df.dtypes.iloc[c], np.number):
                entry = "{:.3f}".format(entry)
            elif np.issubdtype(df.dtypes.iloc[c], np.bool):
                entry = str(entry)
            graphtab.table.cell(r + 1, c).text = entry
            graphtab.table.cell(r + 1, c).text_frame._set_font(*font)


def layout_table_slide(slide, title_text, protopath, protomtgpath, tabledf):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.text_frame._set_font("Candara", 36, False, False)
    pic1 = slide.shapes.add_picture(protomtgpath, Inches(0.0), Inches(0.83333), )
    pic2 = slide.shapes.add_picture(protopath, Inches(0.0), Inches(0.83333), )
    graphtab = slide.shapes.add_table(tabledf.shape[0] + 1, tabledf.shape[1], 0, 0, Inches(5.66), Inches(3.3))
    dataframe2pptxtable(tabledf, graphtab, ("Candara", 12, False, False))
    for k, v in {'height': 1.45, 'width': 6.198, 'top': 0.399, 'left': 6.218}.items():
        setattr(tf, k, Inches(v))
    for k,v in {'height': 4.245, 'width': 10.588, 'top': 3.25, 'left': 2.745}.items():
        setattr(pic1, k, Inches(v))
    for k,v in {'height': 2.17, 'width': 2.17, 'top': 4.12, 'left': 0.00}.items():
        setattr(pic2, k, Inches(v))


def layout_table_slide_all(slide, title_text, protopath, protomtgpath, evolfigpath, maniffigpath, tabledf):
    """
    Layout a slide with a title and two figures.
    """
    tf = slide.shapes.title
    tf.text = title_text
    tf.text_frame._set_font("Candara", 24, False, False)
    pic1 = slide.shapes.add_picture(protomtgpath, Inches(0.0), Inches(0.83333), )
    pic2 = slide.shapes.add_picture(protopath, Inches(0.0), Inches(0.83333), )
    pic3 = slide.shapes.add_picture(evolfigpath, Inches(0.0), Inches(0.83333), )
    pic4 = slide.shapes.add_picture(maniffigpath, Inches(0.0), Inches(0.83333), )
    graphtab = slide.shapes.add_table(tabledf.shape[0] + 1, tabledf.shape[1], 0, 0, Inches(5.66), Inches(3.3))
    dataframe2pptxtable(tabledf, graphtab, ("Candara", 12, False, False))
    for k, v in {'height': 1.45, 'width': 2.745, 'top': 5.677, 'left': 0.0}.items():
        setattr(tf, k, Inches(v))
    for k,v in {'height': 4.245, 'width': 10.588, 'top': 3.25, 'left': 2.745}.items():
        setattr(pic1, k, Inches(v))
    for k,v in {'height': 2.06, 'width': 2.06, 'top': 3.3, 'left': 0.379}.items():
        setattr(pic2, k, Inches(v))
    for k,v in {'height': 3.25, 'width': 3.55, 'top': 0, 'left': 6.572}.items():
        setattr(pic3, k, Inches(v))
    for k,v in {'height': 3.25, 'width': 3.25, 'top': 0, 'left': 10.126}.items():
        setattr(pic4, k, Inches(v))

outdir = join(rootdir, "summary\per_exp_best")
prs = Presentation()
# 16:9 wide screen layout
prs.slide_width = Inches(13.33333)
prs.slide_height = Inches(7.5)
blank_slide_layout = prs.slide_layouts[5]

for Animal in ["Alfa", "Beto"]:
    EStats = loadmat(join(mat_path, Animal + "_Evol_stats.mat"),
                     struct_as_record=False, squeeze_me=True, chars_as_strings=True)['EStats']
    for Expi in range(1, 47):
        if Animal == "Beto" and Expi == 46: continue
        pref_chan = EStats[Expi - 1].evol.pref_chan
        imgpos = EStats[Expi - 1].evol.imgpos
        imgsize = EStats[Expi - 1].evol.imgsize
        area = chan2area(pref_chan)
        title_str = f"{Animal} Exp{Expi:02d} {area} Pref Ch{pref_chan}\npos: {imgpos} {imgsize} deg"
        df_exp_excerpt = pd.read_csv(join(outdir, f"{Animal}_{Expi:02d}_best_methods.csv"), index_col=0)
        protopath = join(sumdir, "proto", f"{Animal}_Exp{Expi:02d}_manif_proto.png")
        mtgpath = join(outdir, f"{Animal}_{Expi:02d}_best_methods_proto.png")
        evolfigpath = join(sumdir, "classic_figs", f"{Animal}_Exp{Expi:02d}_Evolution.png")
        maniffigpath = join(sumdir, "classic_figs", f"{Animal}_Exp{Expi:02d}_Manifold.png")
        slide = prs.slides.add_slide(blank_slide_layout)
        # layout_table_slide(slide, title_str, protopath, mtgpath, df_exp_excerpt)
        layout_table_slide_all(slide, title_str, protopath, mtgpath, evolfigpath, maniffigpath, df_exp_excerpt)
#%%
prs.save(join(sumdir, "ManifExp_best_proto_all_pptx_export.pptx"))
#%%
graphtab = slide.shapes.add_table(df_exp_excerpt.shape[0], df_exp_excerpt.shape[1], 0, 0, Inches(3.173), Inches(4.906))
for r in range(df_exp_excerpt.shape[0]):
    for c in range(df_exp_excerpt.shape[1]):
        entry = df_exp_excerpt.iloc[r, c]
        if np.issubdtype(df_exp_excerpt.dtypes.iloc[c], np.float):
            entry = "{:.3f}".format(entry)
        elif np.issubdtype(df_exp_excerpt.dtypes.iloc[c], np.bool):
            entry = str(entry)
        graphtab.table.cell(r, c).text = entry
